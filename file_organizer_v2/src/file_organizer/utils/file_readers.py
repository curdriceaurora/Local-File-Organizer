"""File reading utilities for various file types."""

from pathlib import Path
from typing import Optional, Union

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub
    EBOOKLIB_AVAILABLE = True
except ImportError:
    EBOOKLIB_AVAILABLE = False

from loguru import logger


class FileReadError(Exception):
    """Exception raised when file reading fails."""
    pass


def read_text_file(file_path: str | Path, max_chars: int = 5000) -> str:
    """Read text content from a plain text file.

    Args:
        file_path: Path to text file
        max_chars: Maximum characters to read

    Returns:
        Text content

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)
    try:
        with open(file_path, encoding='utf-8', errors='ignore') as f:
            text = f.read(max_chars)
        logger.debug(f"Read {len(text)} characters from {file_path.name}")
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read text file {file_path}: {e}") from e


def read_docx_file(file_path: str | Path) -> str:
    """Read text content from a .docx file.

    Args:
        file_path: Path to DOCX file

    Returns:
        Extracted text content

    Raises:
        FileReadError: If file cannot be read
        ImportError: If python-docx is not installed
    """
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx is not installed. Install with: pip install python-docx")

    file_path = Path(file_path)
    try:
        doc = docx.Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        text = '\n'.join(paragraphs)
        logger.debug(f"Extracted {len(text)} characters from {file_path.name}")
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read DOCX file {file_path}: {e}") from e


def read_pdf_file(file_path: str | Path, max_pages: int = 5) -> str:
    """Read text content from a PDF file.

    Args:
        file_path: Path to PDF file
        max_pages: Maximum pages to read

    Returns:
        Extracted text content

    Raises:
        FileReadError: If file cannot be read
        ImportError: If PyMuPDF is not installed
    """
    if not PYMUPDF_AVAILABLE:
        raise ImportError("PyMuPDF is not installed. Install with: pip install PyMuPDF")

    file_path = Path(file_path)
    try:
        doc = fitz.open(file_path)
        num_pages = min(max_pages, len(doc))

        pages_text = []
        for page_num in range(num_pages):
            page = doc.load_page(page_num)
            pages_text.append(page.get_text())

        text = '\n'.join(pages_text)
        doc.close()

        logger.debug(
            f"Extracted {len(text)} characters from {num_pages} pages of {file_path.name}"
        )
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read PDF file {file_path}: {e}") from e


def read_spreadsheet_file(file_path: str | Path, max_rows: int = 100) -> str:
    """Read content from Excel or CSV file.

    Args:
        file_path: Path to spreadsheet file
        max_rows: Maximum rows to read

    Returns:
        String representation of data

    Raises:
        FileReadError: If file cannot be read
        ImportError: If pandas is not installed
    """
    if not PANDAS_AVAILABLE:
        raise ImportError("pandas is not installed. Install with: pip install pandas openpyxl")

    file_path = Path(file_path)
    try:
        # Determine file type and read
        if file_path.suffix.lower() == '.csv':
            df = pd.read_csv(file_path, nrows=max_rows)
        elif file_path.suffix.lower() in ('.xlsx', '.xls'):
            df = pd.read_excel(file_path, nrows=max_rows)
        else:
            raise ValueError(f"Unsupported spreadsheet format: {file_path.suffix}")

        # Convert to string, limiting size
        text = df.to_string(max_rows=max_rows)

        logger.debug(
            f"Extracted {len(text)} characters from {len(df)} rows of {file_path.name}"
        )
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read spreadsheet file {file_path}: {e}") from e


def read_presentation_file(file_path: str | Path) -> str:
    """Read text content from PowerPoint file.

    Args:
        file_path: Path to PPT/PPTX file

    Returns:
        Extracted text from all slides

    Raises:
        FileReadError: If file cannot be read
        ImportError: If python-pptx is not installed
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")

    file_path = Path(file_path)
    try:
        prs = Presentation(file_path)

        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

            if slide_content:
                slides_text.append(f"Slide {slide_num}: " + " | ".join(slide_content))

        text = '\n'.join(slides_text)
        logger.debug(
            f"Extracted {len(text)} characters from {len(slides_text)} slides of {file_path.name}"
        )
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read presentation file {file_path}: {e}") from e


def read_ebook_file(file_path: str | Path, max_chars: int = 10000) -> str:
    """Read text content from ebook file (EPUB only for now).

    Args:
        file_path: Path to ebook file
        max_chars: Maximum characters to extract

    Returns:
        Extracted text content

    Raises:
        FileReadError: If file cannot be read
        ImportError: If ebooklib is not installed
    """
    if not EBOOKLIB_AVAILABLE:
        raise ImportError("ebooklib is not installed. Install with: pip install ebooklib")

    file_path = Path(file_path)

    # Only support EPUB for now
    if file_path.suffix.lower() != '.epub':
        raise ValueError(f"Unsupported ebook format: {file_path.suffix}. Only .epub supported.")

    try:
        book = epub.read_epub(file_path)

        text_parts = []
        total_chars = 0

        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                content = item.get_content().decode('utf-8', errors='ignore')
                # Basic HTML stripping (simple approach)
                import re
                content = re.sub(r'<[^>]+>', ' ', content)
                content = re.sub(r'\s+', ' ', content).strip()

                if content:
                    text_parts.append(content)
                    total_chars += len(content)

                    if total_chars >= max_chars:
                        break

        text = ' '.join(text_parts)[:max_chars]

        logger.debug(f"Extracted {len(text)} characters from ebook {file_path.name}")
        return text
    except Exception as e:
        raise FileReadError(f"Failed to read ebook file {file_path}: {e}") from e


def read_file(file_path: str | Path, **kwargs) -> str | None:
    """Read content from any supported file type.

    Auto-detects file type and uses appropriate reader.

    Args:
        file_path: Path to file
        **kwargs: Additional arguments passed to specific readers

    Returns:
        Extracted text content, or None if unsupported

    Raises:
        FileReadError: If file cannot be read
    """
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    readers = {
        ('.txt', '.md'): read_text_file,
        ('.docx',): read_docx_file,  # Note: .doc (old binary format) is NOT supported
        ('.pdf',): read_pdf_file,
        ('.csv', '.xlsx', '.xls'): read_spreadsheet_file,
        ('.ppt', '.pptx'): read_presentation_file,
        ('.epub',): read_ebook_file,
    }

    for extensions, reader in readers.items():
        if ext in extensions:
            try:
                return reader(file_path, **kwargs)
            except Exception as e:
                logger.error(f"Error reading {file_path.name}: {e}")
                raise

    logger.warning(f"Unsupported file type: {ext}")
    return None
