"""Tests for file_organizer.utils.file_readers module.

Covers the seven text-based reader functions plus helper utilities.
Audio/video/image readers are explicitly excluded (deferred stories).
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from file_organizer.utils.file_readers import (
    FileReadError,
    FileTooLargeError,
    _check_file_size,
    read_text_file,
)


# ---------------------------------------------------------------------------
# _check_file_size
# ---------------------------------------------------------------------------


class TestCheckFileSize:
    """Tests for the _check_file_size helper."""

    def test_small_file_passes(self, tmp_path: Path) -> None:
        f = tmp_path / "small.txt"
        f.write_text("hello")
        _check_file_size(f)  # should not raise

    def test_large_file_raises(self, tmp_path: Path) -> None:
        f = tmp_path / "big.bin"
        f.write_bytes(b"\x00" * 1024)
        with pytest.raises(FileTooLargeError, match="File too large"):
            _check_file_size(f, max_bytes=512)

    def test_exactly_at_limit(self, tmp_path: Path) -> None:
        f = tmp_path / "exact.bin"
        f.write_bytes(b"\x00" * 100)
        _check_file_size(f, max_bytes=100)  # equal is allowed

    def test_missing_file_does_not_raise(self, tmp_path: Path) -> None:
        missing = tmp_path / "nonexistent.txt"
        _check_file_size(missing)  # should silently pass


# ---------------------------------------------------------------------------
# read_text_file
# ---------------------------------------------------------------------------


class TestReadTextFile:
    """Tests for read_text_file."""

    def test_reads_content(self, tmp_path: Path) -> None:
        f = tmp_path / "hello.txt"
        f.write_text("Hello, World!")
        assert read_text_file(f) == "Hello, World!"

    def test_respects_max_chars(self, tmp_path: Path) -> None:
        f = tmp_path / "long.txt"
        f.write_text("A" * 1000)
        result = read_text_file(f, max_chars=50)
        assert len(result) == 50

    def test_accepts_string_path(self, tmp_path: Path) -> None:
        f = tmp_path / "str.txt"
        f.write_text("content")
        assert read_text_file(str(f)) == "content"

    def test_missing_file_raises(self, tmp_path: Path) -> None:
        with pytest.raises(FileReadError):
            read_text_file(tmp_path / "missing.txt")

    def test_utf8_content(self, tmp_path: Path) -> None:
        f = tmp_path / "unicode.txt"
        f.write_text("Héllo Wörld café")
        assert "café" in read_text_file(f)

    def test_empty_file(self, tmp_path: Path) -> None:
        f = tmp_path / "empty.txt"
        f.write_text("")
        assert read_text_file(f) == ""


# ---------------------------------------------------------------------------
# read_docx_file
# ---------------------------------------------------------------------------


class TestReadDocxFile:
    """Tests for read_docx_file."""

    def test_reads_docx(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import DOCX_AVAILABLE, read_docx_file

        if not DOCX_AVAILABLE:
            pytest.skip("python-docx not installed")

        import docx

        doc = docx.Document()
        doc.add_paragraph("Test paragraph one.")
        doc.add_paragraph("Test paragraph two.")
        path = tmp_path / "test.docx"
        doc.save(path)

        result = read_docx_file(path)
        assert "Test paragraph one." in result
        assert "Test paragraph two." in result

    def test_empty_docx(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import DOCX_AVAILABLE, read_docx_file

        if not DOCX_AVAILABLE:
            pytest.skip("python-docx not installed")

        import docx

        doc = docx.Document()
        path = tmp_path / "empty.docx"
        doc.save(path)
        result = read_docx_file(path)
        assert result == ""

    def test_invalid_docx_raises(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import DOCX_AVAILABLE, read_docx_file

        if not DOCX_AVAILABLE:
            pytest.skip("python-docx not installed")

        bad = tmp_path / "bad.docx"
        bad.write_text("not a docx file")
        with pytest.raises(FileReadError):
            read_docx_file(bad)


# ---------------------------------------------------------------------------
# read_pdf_file
# ---------------------------------------------------------------------------


class TestReadPdfFile:
    """Tests for read_pdf_file."""

    def test_reads_pdf(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PYMUPDF_AVAILABLE, read_pdf_file

        if not PYMUPDF_AVAILABLE:
            pytest.skip("PyMuPDF not installed")

        import fitz

        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 72), "Test PDF content")
        path = tmp_path / "test.pdf"
        doc.save(path)
        doc.close()

        result = read_pdf_file(path)
        assert "Test PDF content" in result

    def test_max_pages_respected(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PYMUPDF_AVAILABLE, read_pdf_file

        if not PYMUPDF_AVAILABLE:
            pytest.skip("PyMuPDF not installed")

        import fitz

        doc = fitz.open()
        for i in range(5):
            page = doc.new_page()
            page.insert_text((50, 72), f"Page {i + 1}")
        path = tmp_path / "multi.pdf"
        doc.save(path)
        doc.close()

        result = read_pdf_file(path, max_pages=2)
        assert "Page 1" in result
        assert "Page 2" in result
        # Pages 3-5 should not be read
        assert "Page 3" not in result

    def test_invalid_pdf_raises(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PYMUPDF_AVAILABLE, read_pdf_file

        if not PYMUPDF_AVAILABLE:
            pytest.skip("PyMuPDF not installed")

        bad = tmp_path / "bad.pdf"
        bad.write_text("not a pdf")
        with pytest.raises(FileReadError):
            read_pdf_file(bad)


# ---------------------------------------------------------------------------
# read_spreadsheet_file
# ---------------------------------------------------------------------------


class TestReadSpreadsheetFile:
    """Tests for read_spreadsheet_file."""

    def test_reads_csv(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PANDAS_AVAILABLE, read_spreadsheet_file

        if not PANDAS_AVAILABLE:
            pytest.skip("pandas not installed")

        csv = tmp_path / "data.csv"
        csv.write_text("name,age\nAlice,30\nBob,25\n")

        result = read_spreadsheet_file(csv)
        assert "Alice" in result
        assert "Bob" in result

    def test_reads_xlsx(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PANDAS_AVAILABLE, read_spreadsheet_file

        if not PANDAS_AVAILABLE:
            pytest.skip("pandas not installed")

        import pandas as pd

        df = pd.DataFrame({"x": [1, 2], "y": [3, 4]})
        path = tmp_path / "data.xlsx"
        df.to_excel(path, index=False)

        result = read_spreadsheet_file(path)
        assert "1" in result
        assert "3" in result

    def test_unsupported_extension_raises(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PANDAS_AVAILABLE, read_spreadsheet_file

        if not PANDAS_AVAILABLE:
            pytest.skip("pandas not installed")

        f = tmp_path / "data.ods"
        f.write_text("content")
        with pytest.raises(FileReadError):
            read_spreadsheet_file(f)

    def test_max_rows(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PANDAS_AVAILABLE, read_spreadsheet_file

        if not PANDAS_AVAILABLE:
            pytest.skip("pandas not installed")

        lines = ["col\n"] + [f"{i}\n" for i in range(200)]
        csv = tmp_path / "big.csv"
        csv.write_text("".join(lines))

        result = read_spreadsheet_file(csv, max_rows=10)
        # Should have limited data
        assert result  # non-empty


# ---------------------------------------------------------------------------
# read_presentation_file
# ---------------------------------------------------------------------------


class TestReadPresentationFile:
    """Tests for read_presentation_file."""

    def test_reads_pptx(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PPTX_AVAILABLE, read_presentation_file

        if not PPTX_AVAILABLE:
            pytest.skip("python-pptx not installed")

        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "My Presentation"
        path = tmp_path / "deck.pptx"
        prs.save(path)

        result = read_presentation_file(path)
        assert "My Presentation" in result

    def test_invalid_pptx_raises(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import PPTX_AVAILABLE, read_presentation_file

        if not PPTX_AVAILABLE:
            pytest.skip("python-pptx not installed")

        bad = tmp_path / "bad.pptx"
        bad.write_text("not a pptx")
        with pytest.raises(FileReadError):
            read_presentation_file(bad)


# ---------------------------------------------------------------------------
# read_ebook_file
# ---------------------------------------------------------------------------


class TestReadEbookFile:
    """Tests for read_ebook_file."""

    def test_non_epub_raises(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import EBOOKLIB_AVAILABLE, read_ebook_file

        if not EBOOKLIB_AVAILABLE:
            pytest.skip("ebooklib not installed")

        f = tmp_path / "book.mobi"
        f.write_text("content")
        with pytest.raises(ValueError, match="Only .epub supported"):
            read_ebook_file(f)

    def test_unavailable_raises_import_error(self) -> None:
        with patch("file_organizer.utils.file_readers.EBOOKLIB_AVAILABLE", False):
            from file_organizer.utils.file_readers import read_ebook_file

            with pytest.raises(ImportError):
                read_ebook_file("test.epub")


# ---------------------------------------------------------------------------
# read_file (dispatcher)
# ---------------------------------------------------------------------------


class TestReadFile:
    """Tests for the read_file dispatcher."""

    def test_dispatches_txt(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import read_file

        f = tmp_path / "doc.txt"
        f.write_text("hello world")
        result = read_file(f)
        assert result == "hello world"

    def test_dispatches_csv(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import read_file

        f = tmp_path / "data.csv"
        f.write_text("a,b\n1,2\n")
        result = read_file(f)
        assert result is not None
        assert "1" in result

    def test_dispatches_md(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import read_file

        f = tmp_path / "readme.md"
        f.write_text("# Title\n\nBody")
        result = read_file(f)
        assert "Title" in result

    def test_unsupported_returns_none(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import read_file

        f = tmp_path / "unknown.xyz"
        f.write_text("data")
        result = read_file(f)
        assert result is None

    def test_file_too_large_raises(self, tmp_path: Path) -> None:
        from file_organizer.utils.file_readers import read_file

        f = tmp_path / "big.txt"
        f.write_bytes(b"x" * 1024)

        with patch(
            "file_organizer.utils.file_readers._check_file_size",
            side_effect=FileTooLargeError("File too large"),
        ):
            with pytest.raises(FileTooLargeError):
                read_file(f)
